from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# Skapa en ny presentation
prs = Presentation()

# 🔹 1. Lägg till en titelslide
slide_layout = prs.slide_layouts[0]  
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Welcome to Python-PPTX"
slide.placeholders[1].text = "Creating PowerPoint Presentations with Python"

# 🔹 2. Lägg till innehållsslides
slides_content = [
    ("Content Slide", "This is an example of a content slide."),
    ("Second Slide", "This is another example of a content slide."),
    ("Third Slide", "This is yet another example of a content slide."),
    ("Forth Slide", "This is yet another example of a content slide.")
]

for title_text, content_text in slides_content:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = content_text

# 🔹 3. Lägg till en tabell
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Tom layout
title = slide.shapes.title
title.text = "Example Table"

rows, cols = 4, 3
table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(6), Inches(2)).table

# Fyll tabellen med data
table.cell(0, 0).text = "Category"
table.cell(0, 1).text = "Value 1"
table.cell(0, 2).text = "Value 2"

data = [("A", 10, 20), ("B", 15, 25), ("C", 12, 22)]
for i, (cat, v1, v2) in enumerate(data, start=1):
    table.cell(i, 0).text = cat
    table.cell(i, 1).text = str(v1)
    table.cell(i, 2).text = str(v2)

# 🔹 4. Lägg till ett stapeldiagram
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Tom layout
title = slide.shapes.title
title.text = "Bar Chart Example"

chart_data = CategoryChartData()
chart_data.categories = ["A", "B", "C"]
chart_data.add_series("Series 1", (10, 15, 12))
chart_data.add_series("Series 2", (20, 25, 22))

x, y, cx, cy = Inches(1), Inches(1.5), Inches(6), Inches(3)
chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart

# 🔹 5. Lägg till ett linjediagram
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Tom layout
title = slide.shapes.title
title.text = "Line Chart Example"

chart_data_line = CategoryChartData()
chart_data_line.categories = ["Jan", "Feb", "Mar", "Apr"]
chart_data_line.add_series("Sales 2023", (150, 200, 250, 300))
chart_data_line.add_series("Sales 2024", (180, 220, 270, 320))

chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data_line).chart

# 🔹 6. Lägg till en bild
image_path = "example_image.png"  # Byt ut mot din bildfil
try:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Tom layout
    title = slide.shapes.title
    title.text = "Image Example"

    img_x, img_y, img_width, img_height = Inches(1), Inches(1.5), Inches(5), Inches(3)
    slide.shapes.add_picture(image_path, img_x, img_y, img_width, img_height)
except FileNotFoundError:
    print("Bildfilen saknas! Lägg en bild i samma mapp som skriptet.")

# 🔹 7. Spara presentationen
prs.save("example_presentation2.pptx")
print("Presentation skapad: example_presentation2.pptx")
